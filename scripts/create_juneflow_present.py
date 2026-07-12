from pathlib import Path
import re, textwrap, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

ROOT = Path('/Users/innovera/Documents/juneflow')
OUT = ROOT / 'Juneflow_Investor_Function_Deck_2026-07-08.pptx'
IMG = ROOT / 'tests/visual/reference/shots'
FUNC = ROOT / 'docs/handoff/FUNCTIONS.md'
INV = ROOT / 'docs/handoff/FUNCTIONS-INVENTORY.md'
NAV = ROOT / 'docs/extract/NAV-ROUTES.md'
PTYPE = ROOT / 'docs/extract/PROJECT-TYPES.md'

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# theme colors
NAVY = RGBColor(11, 42, 74)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(14, 165, 233)
GOLD = RGBColor(245, 158, 11)
GREEN = RGBColor(16, 185, 129)
RED = RGBColor(239, 68, 68)
PURPLE = RGBColor(109, 40, 217)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
LIGHT = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(203, 213, 225)

FONT = 'Aptos'
THAI_FONT = 'Arial'

def set_font(run, size=18, bold=False, color=INK):
    run.font.name = THAI_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 250, 252)
    # top accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.13))
    shape.fill.solid(); shape.fill.fore_color.rgb = NAVY; shape.line.fill.background()

def title(slide, text, subtitle=None, y=0.38):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(14.8), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.name = THAI_FONT; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = NAVY
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.63), Inches(y+0.62), Inches(13.8), Inches(0.36))
        p = st.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = THAI_FONT; p.font.size = Pt(13); p.font.color.rgb = MUTED

def footer(slide, n=None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(8.55), Inches(14.8), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.text = 'Juneflow · Construction ERP + Subscription SaaS · Source: docs/handoff + docs/extract'
    p.font.name = THAI_FONT; p.font.size = Pt(8.5); p.font.color.rgb = RGBColor(100,116,139)
    if n:
        p.alignment = PP_ALIGN.RIGHT

def add_card(slide, x, y, w, h, heading, body, color=BLUE, icon=None, heading_size=15, body_size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = BORDER; sh.line.width = Pt(1)
    # accent
    ac = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    ac.fill.solid(); ac.fill.fore_color.rgb = color; ac.line.fill.background()
    tf = sh.text_frame; tf.clear(); tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.16); tf.margin_top = Inches(0.12); tf.word_wrap=True
    p = tf.paragraphs[0]
    p.text = (icon + ' ' if icon else '') + heading
    p.font.name = THAI_FONT; p.font.size = Pt(heading_size); p.font.bold = True; p.font.color.rgb = color
    p2 = tf.add_paragraph(); p2.text = body; p2.space_before=Pt(4)
    p2.font.name = THAI_FONT; p2.font.size = Pt(body_size); p2.font.color.rgb = MUTED
    return sh

def add_bullets(slide, x, y, w, h, bullets, size=15, color=INK, gap=0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap=True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = THAI_FONT; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(gap)
    return tb

def add_kpi(slide, x, y, num, label, color=BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1.15))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    tf = sh.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.text=num; p.alignment=PP_ALIGN.CENTER
    p.font.name=THAI_FONT; p.font.size=Pt(28); p.font.bold=True; p.font.color.rgb=WHITE
    p2=tf.add_paragraph(); p2.text=label; p2.alignment=PP_ALIGN.CENTER
    p2.font.name=THAI_FONT; p2.font.size=Pt(10.5); p2.font.color.rgb=WHITE

def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    return s

def add_image(slide, x, y, w, h, path, caption=None):
    if not Path(path).exists():
        return
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if caption:
        cap = slide.shapes.add_textbox(Inches(x), Inches(y+h+0.05), Inches(w), Inches(0.25))
        p=cap.text_frame.paragraphs[0]; p.text=caption; p.alignment=PP_ALIGN.CENTER
        p.font.name=THAI_FONT; p.font.size=Pt(9); p.font.color.rgb=MUTED

def arrow(slide, x1,y1,x2,y2,color=BLUE):
    line=slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb=color; line.line.width=Pt(2.2)
    return line

# 1 Cover
s=blank()
shape=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
shape.fill.solid(); shape.fill.fore_color.rgb=NAVY; shape.line.fill.background()
# subtle cards
for i,(x,y,c) in enumerate([(10.5,1.2,CYAN),(11.4,2.1,GOLD),(9.8,3.2,GREEN),(12.2,4.4,PURPLE)]):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.0), Inches(1.15))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.fill.transparency=18; sh.line.fill.background()

tb=s.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(9.5), Inches(1.0)); p=tb.text_frame.paragraphs[0]
p.text='Juneflow'; p.font.name=THAI_FONT; p.font.size=Pt(56); p.font.bold=True; p.font.color.rgb=WHITE
tb=s.shapes.add_textbox(Inches(0.85), Inches(2.05), Inches(9.5), Inches(0.75)); p=tb.text_frame.paragraphs[0]
p.text='Construction ERP + Subscription SaaS'; p.font.name=THAI_FONT; p.font.size=Pt(26); p.font.color.rgb=RGBColor(186,230,253)
add_bullets(s,0.9,3.25,8.3,2.2,[
    'ระบบเดียวสำหรับควบคุมต้นทุน งานก่อสร้าง ผู้รับเหมา การเงิน บัญชี ฝ่ายขาย และบริการหลังขาย',
    'ออกแบบเป็น Multi-tenant SaaS: แพ็กเกจ S/M/L/Full, จำกัดเมนู/โควต้า/AI credit ตามแผน',
    'สรุปฟังก์ชันจาก prototype และ handoff documents เพื่อใช้เสนอนักลงทุน'
],17,WHITE,6)
add_kpi(s,0.9,6.55,'44+','เมนูหลัก',CYAN); add_kpi(s,3.95,6.55,'100+','หน้าจอ/route',GOLD); add_kpi(s,7.0,6.55,'719','declared functions',GREEN); add_kpi(s,10.05,6.55,'4','project types',PURPLE)

# 2 Executive summary
s=blank(); title(s,'สรุปใน 1 หน้า','Juneflow คือ ERP เฉพาะทางที่เอา construction workflow + finance + SaaS platform มาอยู่ในระบบเดียว')
add_card(s,0.7,1.45,4.6,2.0,'Pain','ธุรกิจก่อสร้าง/อสังหาฯ ใช้ Excel, แชต, เอกสารกระดาษ และหลายระบบแยกกัน ทำให้ต้นทุน-งบ-งานจริง reconcile ยาก',RED,'01')
add_card(s,5.7,1.45,4.6,2.0,'Product','ตั้งแต่ land → BOQ/BOM → PR/PO/WO/GR → ผู้รับเหมา → GL/AP/AR/Tax → Sales/After-sales → Mobile/LINE',BLUE,'02')
add_card(s,10.7,1.45,4.6,2.0,'Business','ขายแบบ subscription package พร้อม seat, project, storage, AI quota และ admin console สำหรับจัดการ tenant',GREEN,'03')
add_bullets(s,1.0,4.2,14.0,2.0,[
    'จุดเด่นสำหรับนักลงทุน: vertical SaaS ที่มี workflow ลึก ไม่ใช่ generic ERP',
    'รองรับ 4 market entry: อสังหา, Solar/EPC, Civil contractor, Service/maintenance',
    'ข้อมูลทุกโมดูลไหลเข้าศูนย์กลาง: AuditLog, DMS, approval matrix, GL posting, reports hub'
],19,INK,8)
add_image(s,10.3,5.8,4.7,2.2,IMG/'pm-dash.png','ตัวอย่างหน้าจอ dashboard / PM')
footer(s)

# 3 Product map diagram
s=blank(); title(s,'ภาพรวมระบบ: From Project to Profit','แสดงระบบเป็น value chain เดียวตั้งแต่เริ่มโครงการจนถึงรับรู้รายได้')
stages=[('Land / Project','จัดหาที่ดิน\nProject hierarchy'),('Budget','BOQ/BOM\nAI QTO'),('Procure','PR/PO/WO\nGR'),('Build','ผู้รับเหมา\nTimeline / Inventory'),('Finance','AP/AR/GL\nTax / Bank'),('Sell & Serve','CRM / โอน\nAfter-sales / LINE')]
colors=[PURPLE,BLUE,CYAN,GREEN,GOLD,RED]
for i,(h,b) in enumerate(stages):
    x=0.65+i*2.55
    add_card(s,x,2.0,2.25,1.55,h,b,colors[i],None,13,10)
    if i<5: arrow(s,x+2.25,2.77,x+2.52,2.77,colors[i])
add_card(s,1.0,4.7,3.1,1.55,'Platform Layer','Multi-tenant, package builder, quotas, user/role, menu gating',NAVY,'SaaS')
add_card(s,4.45,4.7,3.1,1.55,'Control Layer','Approval matrix, audit log, DMS, notifications, reports hub',BLUE,'CTRL')
add_card(s,7.9,4.7,3.1,1.55,'Data Layer','Master data, project type, cost center, document numbering',GREEN,'DATA')
add_card(s,11.35,4.7,3.1,1.55,'Channel Layer','Web, mobile approval, field app, LINE OA / LIFF',GOLD,'MOB')
add_image(s,0.85,6.8,6.6,1.25,IMG/'land-pipeline.png','ตัวอย่าง land pipeline')
footer(s)

# 4 Market/project types
s=blank(); title(s,'รองรับ 4 ประเภทโครงการ','แต่ละ project type เปิดโมดูลตามบริบท ทำให้ขายได้หลาย vertical จาก product core เดียว')
ptype_cards=[('อสังหาริมทรัพย์','Real Estate','บ้านจัดสรร/คอนโด: ก่อสร้าง → ขาย → หลังขาย',BLUE),('Solar / EPC','Energy','EPC, PPA, O&M, ROI, Permit, Warranty',GOLD),('ก่อสร้างโยธา','Civil','BOQ, จัดซื้อ, ผู้รับเหมา, timeline, inventory',GREEN),('บริการทั่วไป','Service','งบประมาณ, procurement, PM, timeline',PURPLE)]
for i,(th,en,body,c) in enumerate(ptype_cards):
    add_card(s,0.8+(i%2)*7.25,1.55+(i//2)*2.05,6.65,1.55,th+' · '+en,body,c)
# matrix
add_bullets(s,1.0,6.0,14.0,1.3,[
    'Module gating: land, BOQ, procurement, subcon, inventory, PM, sales, LINE OA, O&M, PPA, ROI, permit, warranty',
    'Hierarchy ปรับตามประเภท เช่น Real Estate = โครงการ→เฟส→บล็อก→ยูนิต→Model / Solar = Site→Array→String→Inverter'
],16,INK,6)
footer(s)

# 5 SaaS model
s=blank(); title(s,'SaaS Platform & Subscription','Juneflow มี platform layer สำหรับขายเป็นแพ็กเกจ ไม่ใช่แค่ software project รายครั้ง')
add_card(s,0.8,1.45,4.4,1.7,'Package Builder','กำหนดแพ็กเกจ S/M/L/Full, ราคา, จำนวน project/user/storage/AI credit และเมนูที่เปิดใช้',BLUE,'01')
add_card(s,5.8,1.45,4.4,1.7,'Tenant Control','ลูกค้าเห็นเฉพาะเมนูตาม package + project type; dashboard และ billing ใช้งานเสมอ',GREEN,'02')
add_card(s,10.8,1.45,4.4,1.7,'Admin Console','MRR, subscribers, invoices, package/seat change, suspend, notify, reset password',GOLD,'03')
add_card(s,1.4,4.3,4.0,1.5,'AI Quota','ตัด credit เมื่อใช้ AI QTO · หมดแล้วขึ้น upgrade modal',PURPLE,'AI')
add_card(s,6.0,4.3,4.0,1.5,'Upsell Path','เมนู/โควต้าที่ถูก block ส่งกลับหน้า upgrade package',CYAN,'UP')
add_card(s,10.6,4.3,4.0,1.5,'Billing','ประวัติบิล ใบเสร็จ และสถานะค้างชำระ',RED,'฿')
footer(s)

# 6 Core functions overview
s=blank(); title(s,'Function Coverage — แผนที่โมดูลหลัก','หน้าหลักอ่านง่าย แต่ยังครอบคลุมฟังก์ชันสำคัญทั้งหมด')
modules=[('Master Data','Company/Org, Project Type, Project/Phase/Unit, Vendor/Customer, Cost Center, Doc Numbering, Users/Roles',NAVY),('BOQ & AI QTO','BOQ list/editor/approval/reports, BOM template, AI CAD/BIM takeoff, 3D review, create BOQ',BLUE),('Procurement','PR, PO, WO, GR, vendor compare, variation order, milestone, return/defect',CYAN),('Subcontractor','Contract register, progress claim, inspection, retention, DMS handover, acceptance center',GREEN),('Finance & Accounting','AP, AR, GL, Bank, Tax, FA, Labor, OPEX, RevRec/WIP, project P&L, cashflow',GOLD),('Sales & Service','CRM, quote/booking/contract, down payment, loan/transfer, after-sales ticket',RED),('PM / CMMS','Maintenance contract, schedule, assets, work order, checklist, photo/video, signature, LINE cert',PURPLE),('Mobile / LINE / DMS','Mobile approvals/field/store/foreman, LINE OA/LIFF, DMS, notifications, reports, audit',RGBColor(14,116,144))]
for i,(h,b,c) in enumerate(modules):
    add_card(s,0.65+(i%2)*7.6,1.35+(i//2)*1.62,7.05,1.18,h,b,c,None,13,9.7)
footer(s)

# 7 BOQ AI
s=blank(); title(s,'BOQ, BOM & AI Quantity Takeoff','หัวใจของ construction ERP: คุมงบตั้งแต่แบบจนเกิดเอกสารจัดซื้อ')
steps=[('Upload','IFC/RVT/DWG/DXF/PDF\nLOD + thumbnail + AI credit'),('Process','parse → detect → classify\nmap → calculate'),('Review','แก้ mapping/qty\nconfidence + 3D/2D viewer'),('Create BOQ','สร้าง BOQ พร้อม traceability\nเข้า approval/reports')]
for i,(h,b) in enumerate(steps):
    add_card(s,0.9+i*3.75,1.75,3.1,1.55,h,b,[BLUE,CYAN,GREEN,GOLD][i],str(i+1),13,9.5)
    if i<3: arrow(s,4.0+i*3.75,2.52,4.55+i*3.75,2.52)
add_bullets(s,0.95,4.05,6.6,2.1,[
    'BOQ Editor: เพิ่ม/แก้/ย้ายหมวด, bulk, สร้าง PR จากรายการที่เลือก',
    'Budget Control: เทียบ budget/used/committed และเตือนเกินงบ',
    'Approval: diff เวอร์ชัน, ladder, comment, approve/reject/request change',
    'Reports: cost summary, M/S/L, variance, EVM, revise history'
],14,INK,5)
add_image(s,8.0,3.7,6.8,3.0,IMG/'01-ptype-modal.png','ตัวอย่าง modal / type-aware UI')
footer(s)

# 8 Procurement + subcon
s=blank(); title(s,'Procurement + Subcontractor + Acceptance','เอกสารจัดซื้อและตรวจรับเชื่อมกับ BOQ, DMS, defect, AP และ dashboard')
flow=[('PR','เลือกจาก BOQ\nแนบไฟล์\napproval matrix'),('PO / WO','เทียบ vendor\nVO งานเพิ่ม-ลด\nmilestone'),('GR / Accept','รับเต็ม/บางส่วน\nรูป + ตีกลับ\ndefect center'),('AP','ตั้งหนี้ 3-way match\nPV + WHT\nจ่ายเงิน')]
for i,(h,b) in enumerate(flow):
    add_card(s,1.0+i*3.65,1.7,3.0,1.6,h,b,[BLUE,CYAN,GREEN,GOLD][i],None,14,10)
    if i<3: arrow(s,4.0+i*3.65,2.5,4.6+i*3.65,2.5)
add_card(s,1.0,4.35,6.2,1.55,'ผู้รับเหมา','ทะเบียนสัญญา → งวดงาน → ส่งมอบเอกสาร/รูป → นัดตรวจ → ผ่านจ่าย/ตีกลับ defect → ตรวจซ้ำ → ปิดงาน',GREEN,'SUB')
add_card(s,7.8,4.35,6.2,1.55,'Acceptance Center','รวมตรวจรับทุกชนิด: GR ตีกลับ, งวดงาน, บ้านลูกค้า; เชื่อม Mobile, Dashboard, Audit, Reports, DMS',RED,'ACC')
footer(s)

# 9 Finance
s=blank(); title(s,'Finance & Accounting Automation','ทุกเอกสารเงินไหลเข้า GL posting และรายงานผู้บริหาร')
finance=[('AP','Billing, PV, CN/DN, Deposit, Retention, Aging'),('AR','Invoice, Tax invoice, RV, Credit note, Aging'),('GL','COA, JV, Posting Inbox, Trial, Statements, RevRec/WIP, Close'),('Bank/Tax/FA','Cheque, Recon, Export, VAT/WHT/e-Tax, Asset, Depreciation, Write-off'),('Labor/OPEX','Attendance, payroll, cost allocation, OPEX budget, multi-year comparison'),('Management','Cashflow, Project P&L, reports hub, audit trail')]
for i,(h,b) in enumerate(finance):
    add_card(s,0.75+(i%3)*5.0,1.35+(i//3)*2.0,4.55,1.45,h,b,[BLUE,CYAN,GOLD,GREEN,PURPLE,RED][i],None,13,9.5)
add_bullets(s,1.0,6.1,13.6,0.9,['กติการวม: เงินทุกใบต้องเข้า GL Posting → JV, mutation ทุกจุดเข้า AuditLog, ไฟล์แนบทุกจุดเข้า DMS'],17,NAVY)
footer(s)

# 10 Sales and channels
s=blank(); title(s,'Sales, After-sales, Mobile & LINE OA','ขยายจาก ERP ก่อสร้างไปถึง customer journey หลังขาย')
add_card(s,0.8,1.45,4.4,1.65,'Sales CRM','Lead pipeline, dashboard, lead detail, report dialog, quote/booking/contract',BLUE,'CRM')
add_card(s,5.8,1.45,4.4,1.65,'ขายยูนิต','Quote → จอง → สัญญา → งวดดาวน์ → สินเชื่อ → โอน',GREEN,'SALE')
add_card(s,10.8,1.45,4.4,1.65,'After-sales','แจ้งซ่อม, ticket detail, priority/status, defect loop, รูป/เอกสาร',RED,'SERV')
add_card(s,0.8,4.0,4.4,1.65,'Mobile','อนุมัติเอกสาร, แจ้งซ่อม, PM technician, store, foreman, GR/defect',PURPLE,'APP')
add_card(s,5.8,4.0,4.4,1.65,'LINE OA / LIFF','บิล, นัดหมาย, แจ้งซ่อม, งวดดาวน์, PM quote/cert/contracts, rich menu',GREEN,'LINE')
add_image(s,10.8,3.72,4.15,2.6,IMG/'pm-manual4.png','ตัวอย่าง mobile/PM flow')
footer(s)

# 11 PM CMMS
s=blank(); title(s,'PM / CMMS: งานบริการและบำรุงรักษา','เปิดตลาดบริการ/maintenance และเชื่อม field workflow กับ LINE')
add_image(s,0.85,1.35,6.8,3.8,IMG/'pm-grouped.png','ตัวอย่าง PM grouped view')
add_card(s,8.0,1.55,3.25,1.4,'Contract','wizard สร้างสัญญา, MA/รายครั้ง, ต่ออายุ, project-first',BLUE)
add_card(s,11.7,1.55,3.25,1.4,'Schedule','auto-gen ใบงานลงปฏิทิน, คลิกวันเพื่อกรองงาน',CYAN)
add_card(s,8.0,3.45,3.25,1.4,'Work Order','check-in GPS, checklist, รูปก่อน/หลัง, signature',GREEN)
add_card(s,11.7,3.45,3.25,1.4,'Customer Channel','ใบรับรองผล → LINE, อะไหล่ → quote → ลูกค้าอนุมัติผ่าน LINE',GOLD)
add_bullets(s,8.0,6.0,6.9,1.0,['เหมาะกับ recurring revenue เพิ่มจาก software: maintenance contracts + after-sales + field operations'],16,INK)
footer(s)

# 12 Land/Solar
s=blank(); title(s,'Land, Solar/EPC และ Vertical Expansion','ใช้ core เดียว แต่เปิดฟีเจอร์เฉพาะ vertical ตาม project type')
add_card(s,0.8,1.45,4.6,1.65,'Land','Pipeline 7 ขั้น, Land Bank, survey/feasibility, DD checklist, purchase/lease contract',PURPLE,'LAND')
add_card(s,5.7,1.45,4.6,1.65,'Solar / EPC','Monitoring inverter, O&M ticket, PPA billing, ROI, permit, warranty',GOLD,'SOLAR')
add_card(s,10.6,1.45,4.6,1.65,'Civil / Service','BOQ/procurement/subcon/timeline สำหรับ civil; service เน้นงบ/แผน/PM',GREEN,'CIV')
add_image(s,0.9,4.0,6.6,3.0,IMG/'03-land-rest.png','Land workflow reference')
add_image(s,8.2,4.0,6.6,3.0,IMG/'02-pm-c-wizard.png','Project-first wizard reference')
footer(s)

# 13 Architecture principles
s=blank(); title(s,'Operational Backbone: สิ่งที่ทำให้ระบบ Scale ได้','กฎกลางที่บังคับทุกโมดูลเพื่อความน่าเชื่อถือของข้อมูล')
backbone=[('AuditLog','ทุก mutation ต้องมี log ใคร-ทำอะไร-เมื่อไหร่'),('DMS','ทุกไฟล์แนบเข้าศูนย์เอกสาร พร้อม link_module'),('Approval','สถานะ draft→pending→approved/rejected ตาม matrix'),('GL Posting','เอกสารเงินทุกใบต้อง post เป็น JV'),('i18n','รองรับ 4 ภาษา th/zh/en/ar-RTL'),('Quota/Middleware','ตรวจ package/menu/AI credit ที่ middleware')]
for i,(h,b) in enumerate(backbone):
    add_card(s,0.8+(i%3)*5.0,1.35+(i//3)*2.2,4.45,1.55,h,b,[NAVY,BLUE,GREEN,GOLD,PURPLE,RED][i],None,14,10)
add_bullets(s,1.0,6.35,14.0,0.8,['Investor message: ความลึกของ workflow + backbone เหล่านี้สร้าง switching cost และทำให้ product แตกต่างจาก SaaS ทั่วไป'],17,NAVY)
footer(s)

# 14 Investor narrative
s=blank(); title(s,'Investor Story ที่ควรเล่า','เน้น “vertical SaaS ที่คุมเงินจริงและ workflow จริง”')
add_card(s,0.8,1.35,4.3,4.7,'1. Problem','ข้อมูลก่อสร้างกระจาย: BOQ อยู่ Excel, จัดซื้ออยู่เอกสาร, ผู้รับเหมาอยู่แชต, บัญชีลงทีหลัง ทำให้คุมต้นทุน/กำไรโครงการช้าและผิดพลาด',RED,None,16,13)
add_card(s,5.85,1.35,4.3,4.7,'2. Solution','Juneflow รวม workflow ทั้งหมดไว้ในระบบเดียว พร้อม approval, audit, DMS, GL posting และ dashboard ผู้บริหารแบบ project-aware',BLUE,None,16,13)
add_card(s,10.9,1.35,4.3,4.7,'3. Scale','ขายแบบ subscription package ให้หลายประเภทโครงการ ใช้ core platform เดียว แต่เปิดโมดูลตาม vertical และเพิ่ม AI/field/mobile เป็น upsell',GREEN,None,16,13)
add_bullets(s,1.0,6.8,13.8,0.7,['ประโยคสั้นบนสไลด์: “From construction chaos to project-level financial control.”'],18,NAVY)
footer(s)

# 15 Suggested deck flow
s=blank(); title(s,'โครงสไลด์สำหรับใช้ Pitch จริง','เวอร์ชันนำเสนอควรใช้ 10–12 หน้า ส่วนรายละเอียดฟังก์ชันเก็บเป็น appendix')
pitch=[('1','Problem','เอกสาร/ต้นทุน/บัญชี/ไซต์งานแยกกัน'),('2','Solution','Construction ERP + SaaS platform'),('3','Product Map','flow ตั้งแต่ land → finance → after-sales'),('4','Core Modules','BOQ, Procurement, Subcon, Finance, Sales, PM'),('5','SaaS Model','แพ็กเกจ + quota + admin console'),('6','Market Entry','Real Estate, Solar/EPC, Civil, Service'),('7','Differentiation','AI QTO, GL automation, Mobile/LINE, Audit/DMS'),('8','Traction / Demo','ใส่ screenshot หรือ demo flow'),('9','Business Model','subscription + AI credit + implementation/support'),('10','Ask / Roadmap','เงินที่ต้องการใช้ทำอะไร')]
for i,(no,h,b) in enumerate(pitch):
    add_card(s,0.85+(i%2)*7.35,1.2+(i//2)*1.35,6.9,0.95,no+'. '+h,b,[BLUE,GREEN,GOLD,PURPLE,RED][i%5],None,12,9.5)
footer(s)

# 16 Demo screenshots
s=blank(); title(s,'ภาพประกอบจาก Prototype / Visual References','ใช้เป็นสไลด์ demo หรือแบ่งใส่ตามโมดูล')
imgs=[('land-pipeline.png','Land pipeline'),('pm-dash.png','PM dashboard'),('pm-contract-step1.png','PM contract wizard'),('02-pm-c-detail.png','PM contract detail')]
for i,(fn,cap) in enumerate(imgs):
    add_image(s,0.75+(i%2)*7.55,1.3+(i//2)*3.25,6.8,2.45,IMG/fn,cap)
footer(s)

# 17 Appendix intro
s=blank(); title(s,'Appendix: Function Coverage','รายละเอียดถัดไปคือ function inventory จาก source document เพื่อยืนยันความครบถ้วน')
add_bullets(s,1.0,1.7,13.8,2.0,[
    'Source: docs/handoff/FUNCTIONS.md — สรุป behavior ราย feature',
    'Source: docs/handoff/FUNCTIONS-INVENTORY.md — สแกน component/function จาก prototype ทั้งหมด',
    'Coverage: 719 declared functions / 78 source files, route/menu จาก docs/extract/NAV-ROUTES.md',
    'แนวทางใช้จริง: presenter ใช้สไลด์หลัก 1–16; appendix เปิดเมื่อตอบคำถามเชิงลึก'
],18,INK,8)
add_card(s,1.0,5.2,13.8,1.5,'หมายเหตุ','เพื่อให้อ่านง่าย สไลด์หลักจะรวม function เป็นกลุ่มธุรกิจ ส่วน appendix เก็บชื่อ function/source file ครบถ้วนแบบ compact',NAVY)
footer(s)

# Appendix from inventory chunks
inventory_text = INV.read_text(encoding='utf-8')
# collect file sections
sections=[]
current=None
for line in inventory_text.splitlines():
    if line.startswith('## '):
        current=line[3:].strip(); sections.append([current,''])
    elif current and line.startswith('- '):
        sections[-1][1]=line[2:].strip()
# chunk sections by approximate char length
chunks=[]; chunk=[]; char=0
for name,body in sections:
    entry=f'{name}: {body}'
    if char+len(entry)>1250 and chunk:
        chunks.append(chunk); chunk=[]; char=0
    chunk.append(entry); char+=len(entry)
if chunk: chunks.append(chunk)
for idx,chunk in enumerate(chunks, start=1):
    s=blank(); title(s,f'Appendix A{idx}: Function Inventory',f'ชื่อไฟล์และ function/component ที่ประกาศจริง — chunk {idx}/{len(chunks)}')
    # two columns text
    col1=chunk[:math.ceil(len(chunk)/2)]; col2=chunk[math.ceil(len(chunk)/2):]
    for col,x in [(col1,0.7),(col2,8.15)]:
        tb=s.shapes.add_textbox(Inches(x), Inches(1.25), Inches(7.05), Inches(7.1))
        tf=tb.text_frame; tf.clear(); tf.word_wrap=True
        for i,entry in enumerate(col):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text='• '+entry
            p.font.name=THAI_FONT; p.font.size=Pt(6.3); p.font.color.rgb=INK
            p.space_after=Pt(1.5)
    footer(s)

# Route appendix summary table-like
s=blank(); title(s,'Appendix B: Route/Menu Coverage','สรุปจำนวน route ต่อกลุ่มเมนูจาก NAV-ROUTES')
route_groups=[('งานหลัก','Land, BOQ, PR/PO/WO/GR, Subcon, Timeline, Inventory, Petty, Acceptance, Labor, PM'),('พลังงาน/EPC','Solar Monitoring, PPA, ROI, Permit, Warranty'),('บัญชี-การเงิน','Alloc, OPEX, GL, AP, AR, Bank, Tax, FA'),('งานขาย-อสังหาฯ','Sales dashboard, CRM, process, down, loan, after-sales'),('ระบบ','Master data, users, reports, DMS, settings, audit, mobile, LINE, sync'),('บัญชีการใช้งาน/Platform','Tenant package/billing + platform admin revenue/subscribers/plans/invoices')]
for i,(h,b) in enumerate(route_groups):
    add_card(s,0.85+(i%2)*7.25,1.35+(i//2)*1.75,6.7,1.25,h,b,[BLUE,GOLD,GREEN,RED,PURPLE,NAVY][i],None,12.5,9)
footer(s)

# Source slide
s=blank(); title(s,'Sources & Next Steps','ไฟล์นี้สร้างจากเอกสารจริงในเครื่อง')
add_bullets(s,0.9,1.45,14.0,2.5,[
    'docs/handoff/FUNCTIONS.md — feature behavior summary',
    'docs/handoff/FUNCTIONS-INVENTORY.md — 719 functions / 78 files inventory',
    'docs/extract/NAV-ROUTES.md — menu/route coverage',
    'docs/extract/PROJECT-TYPES.md — 4 project types + module gating',
    'tests/visual/reference/shots/*.png — ภาพประกอบจาก visual references'
],16,INK,6)
add_card(s,0.9,5.0,6.8,1.5,'Next: Investor-ready version','ถ้าจะใช้ pitch จริง จ่อยแนะนำตัด appendix ออกเป็น backup และเพิ่ม traction, market size, pricing, ask, roadmap',GREEN)
add_card(s,8.1,5.0,6.8,1.5,'Next: Demo version','ถ้าจะใช้ demo product จ่อยแนะนำเพิ่ม screenshot ราย module + highlight flow ที่ investor เข้าใจใน 3 นาที',BLUE)
footer(s)

prs.save(OUT)
print(OUT)
print('slides', len(prs.slides))
print('appendix_chunks', len(chunks))
