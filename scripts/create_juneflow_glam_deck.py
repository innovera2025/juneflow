from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFilter
import math, random, re, zipfile

ROOT = Path('/Users/innovera/Documents/juneflow')
OUT = ROOT / 'Juneflow_Investor_Glam_Deck_2026-07-08.pptx'
ASSETS = ROOT / 'deck_assets_glam'
SHOTS = ROOT / 'tests/visual/reference/shots'
ASSETS.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

W, H = 1920, 1080
NAVY = RGBColor(5, 16, 35)
INK = RGBColor(12, 18, 33)
WHITE = RGBColor(255,255,255)
MUTED = RGBColor(203, 213, 225)
CYAN = RGBColor(34, 211, 238)
BLUE = RGBColor(59, 130, 246)
GOLD = RGBColor(251, 191, 36)
GREEN = RGBColor(52, 211, 153)
PURPLE = RGBColor(168, 85, 247)
RED = RGBColor(248, 113, 113)
THAI = 'Arial'

PALETTE = {
    'navy': (5,16,35), 'navy2': (13,27,55), 'cyan': (34,211,238), 'blue': (59,130,246),
    'gold': (251,191,36), 'green': (52,211,153), 'purple': (168,85,247), 'red': (248,113,113),
    'white': (255,255,255), 'muted': (203,213,225)
}

def pil_bg(name, accent='cyan', constellation=False):
    path = ASSETS / f'{name}.png'
    if path.exists(): return path
    img = Image.new('RGB', (W,H), PALETTE['navy'])
    pix = img.load()
    a = PALETTE[accent]
    for y in range(H):
        for x in range(W):
            dx = (x-W*0.72)/W; dy=(y-H*0.20)/H
            glow = max(0, 1 - (dx*dx*4 + dy*dy*7))
            dx2=(x-W*0.1)/W; dy2=(y-H*0.95)/H
            glow2=max(0,1-(dx2*dx2*8+dy2*dy2*5))
            base=PALETTE['navy']
            r=int(base[0]+glow*a[0]*0.35+glow2*90*0.25)
            g=int(base[1]+glow*a[1]*0.28+glow2*80*0.18)
            b=int(base[2]+glow*a[2]*0.35+glow2*140*0.20)
            pix[x,y]=(min(255,r),min(255,g),min(255,b))
    draw = ImageDraw.Draw(img, 'RGBA')
    random.seed(9)
    for _ in range(70):
        x=random.randint(0,W); y=random.randint(0,H); r=random.choice([1,1,2]); alpha=random.randint(30,100)
        draw.ellipse((x-r,y-r,x+r,y+r), fill=(255,255,255,alpha))
    # light grid / construction lines
    for x in range(-300,W,110):
        draw.line((x,H,x+600,0), fill=(255,255,255,10), width=1)
    for y in range(0,H,90):
        draw.line((0,y,W,y), fill=(255,255,255,7), width=1)
    if constellation:
        pts=[]
        for i in range(44):
            ang = 2*math.pi*i/44
            rad = 260 + 80*math.sin(i*2.7)
            cx = W*0.72 + math.cos(ang)*rad
            cy = H*0.50 + math.sin(ang)*rad*0.72
            pts.append((cx,cy))
        for i,p in enumerate(pts):
            q=pts[(i+5)%len(pts)]
            draw.line((p[0],p[1],q[0],q[1]), fill=(*a,45), width=2)
        for i,(x,y) in enumerate(pts):
            col=[PALETTE['cyan'],PALETTE['gold'],PALETTE['green'],PALETTE['purple']][i%4]
            draw.ellipse((x-7,y-7,x+7,y+7), fill=(*col,220), outline=(255,255,255,130), width=1)
    img.save(path, quality=95)
    return path

def add_bg(slide, name='bg', accent='cyan', constellation=False):
    slide.shapes.add_picture(str(pil_bg(name, accent, constellation)), 0, 0, width=prs.slide_width, height=prs.slide_height)

def tx(slide, text, x,y,w,h, size=24, color=WHITE, bold=False, align='left', alpha=None):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    p=tf.paragraphs[0]; p.text=text
    p.font.name=THAI; p.font.size=Pt(size); p.font.bold=bold; p.font.color.rgb=color
    p.alignment={'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}[align]
    return tb

def title(slide, main, sub=None, dark=False):
    tx(slide, main, .65, .45, 14.8, .55, 24, WHITE if not dark else INK, True)
    if sub: tx(slide, sub, .68, .98, 13.9, .35, 10.8, MUTED if not dark else RGBColor(71,85,105))

def glass(slide, x,y,w,h, fill=(255,255,255), transparency=82, line=(255,255,255), ltrans=72, radius=True):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(*fill); sh.fill.transparency=transparency
    sh.line.color.rgb=RGBColor(*line); sh.line.transparency=ltrans; sh.line.width=Pt(1)
    return sh

def card(slide, x,y,w,h, head, body='', accent=CYAN, num=None, big=False):
    glass(slide,x,y,w,h,transparency=86)
    # accent bar
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(.055), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb=accent; bar.line.fill.background()
    if num:
        tx(slide, num, x+.18,y+.18,.55,.3,10,accent,True)
        hx=x+.82
    else: hx=x+.22
    tx(slide, head, hx,y+.18,w-.45,.38, 15 if not big else 20, WHITE, True)
    if body: tx(slide, body, x+.22,y+.66,w-.45,h-.78, 10.5 if not big else 13, MUTED)

def kpi(slide,x,y,num,label,accent=CYAN):
    glass(slide,x,y,3.0,1.05,transparency=88,line=(255,255,255),ltrans=78)
    tx(slide,num,x+.15,y+.12,2.7,.45,28,accent,True,'center')
    tx(slide,label,x+.18,y+.66,2.65,.22,9.5,MUTED,False,'center')

def pic(slide, x,y,w,h, path, caption=None):
    if not Path(path).exists(): return
    glass(slide,x-.08,y-.08,w+.16,h+.16,transparency=90,line=(34,211,238),ltrans=60)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if caption: tx(slide, caption, x, y+h+.06, w, .25, 8.5, MUTED, False, 'center')

def footer(slide):
    tx(slide,'Juneflow · Investor Pitch Deck · compact premium version',.65,8.58,8,.22,7.5,MUTED)

def blank(accent='cyan', constellation=False):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s,f'bg_{len(prs.slides)}',accent,constellation); return s

# 1 Cover
s=blank('cyan', True)
tx(s,'Juneflow',.75,1.0,8.8,.8,56,WHITE,True)
tx(s,'Construction ERP SaaS\nfor Real Project Control',.82,2.05,7.2,1.05,29,CYAN,True)
tx(s,'จาก BOQ → จัดซื้อ → ผู้รับเหมา → บัญชี → ขาย → หลังขาย\nระบบเดียวที่ทำให้ “ต้นทุนโครงการ” มองเห็นได้จริง',.85,3.45,7.6,.9,16,MUTED)
kpi(s,.85,6.85,'44+','เมนูหลัก',CYAN); kpi(s,4.05,6.85,'100+','screens / routes',GOLD); kpi(s,7.25,6.85,'719','prototype functions',GREEN); kpi(s,10.45,6.85,'4','vertical project types',PURPLE)
footer(s)

# 2 One-liner
s=blank('gold')
title(s,'Investor Summary','ข้อความหลักที่อยากให้นักลงทุนจำได้ใน 20 วินาที')
tx(s,'Juneflow เปลี่ยนงานก่อสร้างที่กระจายอยู่ใน Excel, เอกสาร, แชต และบัญชีปลายทาง\nให้กลายเป็นระบบ ERP SaaS ที่คุมต้นทุนและกำไรได้ระดับโครงการ',1.1,1.8,13.8,1.25,28,WHITE,True,'center')
card(s,1.1,4.0,4.25,1.55,'Pain','ข้อมูลจริงของไซต์งานกับตัวเลขบัญชีไม่เคยอยู่ที่เดียวกัน',RED,'01')
card(s,5.85,4.0,4.25,1.55,'Product','workflow ลึกตั้งแต่ BOQ ถึง GL และ After-sales',BLUE,'02')
card(s,10.6,4.0,4.25,1.55,'Scale','ขายเป็น subscription ตามแพ็กเกจ เมนู โควต้า และ AI credit',GREEN,'03')
footer(s)

# 3 Product galaxy
s=blank('cyan', True)
title(s,'Product Galaxy','ครบทุก function แต่เล่าแบบไม่รก: รวมเป็น 9 กลุ่มธุรกิจ')
modules=[('SaaS Platform','Package / Tenant / Admin'),('Master Data','Company / Project / Role'),('BOQ + AI QTO','Budget / BIM / Approval'),('Procurement','PR / PO / WO / GR'),('Subcontractor','Progress / Accept / Retention'),('Finance','AP / AR / GL / Tax / Bank'),('Sales','CRM / Booking / Transfer'),('PM / CMMS','Contract / Schedule / WO'),('Channels','Mobile / LINE / DMS')]
# center
circle=slide = s
for i,(h,b) in enumerate(modules):
    ang=2*math.pi*i/len(modules); x=6.5+math.cos(ang)*4.6; y=4.05+math.sin(ang)*2.45
    card(s,x,y,3.0,1.0,h,b,[CYAN,BLUE,GOLD,GREEN,PURPLE,RED,CYAN,GREEN,PURPLE][i],str(i+1))
glass(s,6.05,3.35,3.9,1.25,transparency=80,line=(34,211,238),ltrans=30)
tx(s,'Juneflow Core',6.2,3.62,3.6,.35,24,CYAN,True,'center')
tx(s,'Audit · DMS · Approval · GL Posting',6.2,4.08,3.6,.25,10,MUTED,False,'center')
footer(s)

# 4 Value chain
s=blank('green')
title(s,'One Operating System for Construction','เห็น flow เงินและงานตั้งแต่ต้นน้ำถึงปลายน้ำ')
steps=[('Land','หา/ประเมินที่ดิน'),('BOQ','ตั้งงบ + ถอดปริมาณ'),('Procure','PR/PO/WO/GR'),('Build','ผู้รับเหมา + ตรวจรับ'),('Finance','AP/AR/GL/Tax'),('Sell/Serve','ขาย + หลังขาย')]
for i,(h,b) in enumerate(steps):
    x=.75+i*2.55; y=2.3
    glass(s,x,y,2.25,1.35,transparency=84,line=PALETTE['cyan'],ltrans=65)
    tx(s,h,x+.18,y+.22,1.9,.35,16,WHITE,True,'center')
    tx(s,b,x+.15,y+.72,1.95,.32,9.7,MUTED,False,'center')
    if i<5:
        ln=s.shapes.add_connector(1, Inches(x+2.25), Inches(y+.67), Inches(x+2.55), Inches(y+.67))
        ln.line.color.rgb=CYAN; ln.line.width=Pt(2)
pic(s,1.0,5.05,4.4,2.45,SHOTS/'land-pipeline.png','Land pipeline')
pic(s,5.85,5.05,4.4,2.45,SHOTS/'pm-dash.png','PM / dashboard')
pic(s,10.7,5.05,4.4,2.45,SHOTS/'pm-contract-step1.png','Project-first wizard')
footer(s)

# 5 Differentiation
s=blank('purple')
title(s,'Why Juneflow Wins','แตกต่างด้วย workflow ลึก + data backbone ที่ทำให้เปลี่ยนระบบยาก')
card(s,.9,1.55,4.5,2.0,'Vertical Depth','ไม่ใช่ ERP ทั่วไป: มี BOQ, AI QTO, งวดผู้รับเหมา, retention, defect, project P&L',CYAN,'01',True)
card(s,5.75,1.55,4.5,2.0,'Financial Truth','เอกสารเงินทุกใบไหลเข้า GL Posting / JV ทำให้ต้นทุนจริง reconcile ได้',GOLD,'02',True)
card(s,10.6,1.55,4.5,2.0,'Field Channels','Mobile approval, foreman, technician, LINE OA/LIFF ทำให้ข้อมูลจากไซต์เข้า ERP',GREEN,'03',True)
card(s,2.1,4.8,5.2,1.6,'SaaS Monetization','แพ็กเกจ S/M/L/Full, menu gating, quotas, AI credit, upgrade path',PURPLE,'04')
card(s,8.2,4.8,5.2,1.6,'Trust Layer','AuditLog + DMS + approval matrix + i18n 4 ภาษา + middleware quota',BLUE,'05')
footer(s)

# 6 SaaS business model
s=blank('gold')
title(s,'Business Model','ทำเงินได้มากกว่า license: subscription + usage + implementation')
models=[('Subscription','รายเดือน/รายปีตามแพ็กเกจ'),('Seat / Project','คิดตามผู้ใช้หรือจำนวนโครงการ'),('AI Credit','AI QTO / document intelligence'),('Implementation','setup, migration, training'),('Support','premium support / SLA'),('Vertical Upsell','Solar, PM, LINE, advanced finance')]
for i,(h,b) in enumerate(models):
    card(s,.85+(i%3)*5.0,1.55+(i//3)*2.1,4.45,1.5,h,b,[CYAN,GOLD,GREEN,PURPLE,BLUE,RED][i],str(i+1))
tx(s,'Package Builder ใน prototype รองรับ S/M/L/Full + quota project/user/storage/AI แล้ว',1.0,6.65,14,.55,22,WHITE,True,'center')
footer(s)

# 7 Market entry
s=blank('green')
title(s,'Market Entry: 4 Verticals, 1 Core Platform','เริ่มจาก Real Estate/Civil แล้วขยายไป Solar และ Service ได้')
verts=[('Real Estate','ก่อสร้าง → ขายยูนิต → After-sales → LINE ลูกบ้าน',BLUE),('Solar / EPC','Monitoring, O&M, PPA, ROI, Permit, Warranty',GOLD),('Civil Contractor','BOQ, Procurement, Subcontractor, Timeline, Inventory',GREEN),('Service / PM','สัญญาบริการ, ใบงาน, technician, certification',PURPLE)]
for i,(h,b,c) in enumerate(verts):
    card(s,1.0+(i%2)*7.3,1.65+(i//2)*2.05,6.55,1.55,h,b,c,str(i+1),True)
tx(s,'กลยุทธ์: ใช้ core เดียว ลด cost of delivery แต่เปิด module ตาม project type เพื่อขายหลายตลาด',1.0,6.6,14,.55,20,CYAN,True,'center')
footer(s)

# 8 Product demo screenshots
s=blank('cyan')
title(s,'Product Proof: Prototype Screens','ภาพประกอบสำหรับ demo — ใช้โชว์ความเป็นระบบจริง ไม่ใช่แค่ concept')
pic(s,.9,1.55,4.55,2.65,SHOTS/'pm-grouped.png','PM grouped / service operations')
pic(s,5.75,1.55,4.55,2.65,SHOTS/'02-pm-c-detail.png','Contract detail')
pic(s,10.6,1.55,4.55,2.65,SHOTS/'03-land-rest.png','Land workflow')
card(s,1.2,5.55,4.0,1.25,'Demo Story','เริ่มจาก Project → BOQ → PR/PO → ตรวจรับ → GL',CYAN)
card(s,6.0,5.55,4.0,1.25,'Investor Lens','โชว์ว่า data ไหลจาก operation ไป finance จริง',GOLD)
card(s,10.8,5.55,4.0,1.25,'Next Asset','ควรทำ video demo 90 วินาทีต่อจาก deck นี้',GREEN)
footer(s)

# 9 Function coverage compact
s=blank('purple')
title(s,'Function Coverage, Condensed','ยังครบ แต่ไม่ทำให้คนอ่านจมในรายละเอียด')
coverage=[('Core Shell','route, modal, toast, project switcher, i18n, sidebar gating'),('Commercial','subscription admin, tenant package, billing, quotas, AI credit'),('Construction','BOQ/BOM/AI QTO, PR/PO/WO/GR, subcon, timeline, inventory'),('Finance','AP, AR, GL, Bank, Tax, FA, Labor, OPEX, RevRec/WIP'),('Customer','Sales CRM, quote, booking, contract, down, loan, after-sales'),('Channels & Control','Mobile, LINE OA, DMS, reports, audit, settings, sync')]
for i,(h,b) in enumerate(coverage):
    card(s,.8+(i%2)*7.35,1.45+(i//2)*1.75,6.75,1.22,h,b,[CYAN,GOLD,GREEN,BLUE,RED,PURPLE][i],str(i+1))
kpi(s,3.2,7.0,'719','functions scanned',CYAN); kpi(s,6.5,7.0,'78','prototype files',GOLD); kpi(s,9.8,7.0,'110±','routes/screens',GREEN)
footer(s)

# 10 Investor ask placeholder
s=blank('blue')
title(s,'What to Raise For','สไลด์นี้ตั้งใจทำเป็น placeholder ให้พี่เติมตัวเลขจริง')
card(s,.95,1.65,4.25,2.05,'Productization','ทำให้ prototype → production SaaS: contract, auth, billing, audit, deployment',CYAN,'01',True)
card(s,5.85,1.65,4.25,2.05,'Go-to-market','pilot customers, implementation playbook, sales collateral, demo video',GOLD,'02',True)
card(s,10.75,1.65,4.25,2.05,'AI + Data','AI QTO, document intelligence, project cost analytics, benchmark data moat',GREEN,'03',True)
tx(s,'ใส่เพิ่มภายหลัง: เงินที่ขอ · runway · milestones · pilot logos · pricing',1.0,5.55,14,.7,26,WHITE,True,'center')
footer(s)

# 11 Closing
s=blank('cyan', True)
tx(s,'Juneflow',.9,1.25,6,.7,48,WHITE,True)
tx(s,'The financial control layer\nfor construction businesses.',.95,2.25,8.5,1.2,34,CYAN,True)
tx(s,'Vertical SaaS ที่เริ่มจาก workflow จริงของไซต์งาน\nและจบที่ตัวเลขกำไรขาดทุนระดับโครงการ',.98,4.0,7.6,.9,18,MUTED)
card(s,1.0,6.35,4.4,1.15,'Pitch takeaway','คุมงาน + คุมต้นทุน + คุมเงิน ในระบบเดียว',GOLD)
card(s,5.95,6.35,4.4,1.15,'Product moat','workflow ลึก + data backbone + field channels',GREEN)
card(s,10.9,6.35,4.4,1.15,'SaaS upside','package + quota + AI credit + vertical expansion',CYAN)
footer(s)

# 12 Sources small
s=blank('gold')
title(s,'Source Basis','ย่อให้สวย แต่ยังยึดเอกสารจริง')
for i,line in enumerate([
    'docs/handoff/FUNCTIONS.md — feature behavior summary',
    'docs/handoff/FUNCTIONS-INVENTORY.md — 719 functions / 78 files',
    'docs/extract/NAV-ROUTES.md — menu and route coverage',
    'docs/extract/PROJECT-TYPES.md — 4 vertical project types',
    'tests/visual/reference/shots/*.png — prototype screenshots'
]):
    card(s,1.4,1.45+i*1.05,13.2,.72,line,'', [CYAN,BLUE,GOLD,GREEN,PURPLE][i],None)
tx(s,'Deck นี้เป็นเวอร์ชัน “อ่านเร็ว / pitch-ready” รายละเอียด function เต็มอยู่ใน deck เก่าและ source documents',1.4,7.25,13.2,.45,15,MUTED,False,'center')
footer(s)

prs.save(OUT)
print(OUT)
print('slides', len(prs.slides))
print('size_mb', round(OUT.stat().st_size/1024/1024,2))
# verify zip
with zipfile.ZipFile(OUT) as z:
    bad=z.testzip()
print('zip_bad', bad)
